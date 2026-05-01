# document/pptx_watermark.py
"""
PowerPoint文档水印处理器 - 隐藏形状方式

本模块实现 PowerPoint 文档的水印处理：
- 嵌入方式：在文本前添加隐藏标记
- 标记格式：hidden_<水印内容>_<原文>
- 优点：实现简单，不影响幻灯片显示
- 缺点：水印在文本中可见（可通过XML解析隐藏）

技术方案：
1. 使用 python-pptx 打开演示文稿
2. 遍历所有幻灯片和形状
3. 在每个文本框的文本前添加水印标记
4. 提取时解析 XML 查找标记

使用示例:
    >>> handler = PPTXHandler()
    >>> 
    >>> # 嵌入水印
    >>> result = handler.embed(
    ...     "presentation.pptx",
    ...     WatermarkData(content="内部资料"),
    ...     "output.pptx"
    ... )
    >>> 
    >>> # 提取水印
    >>> result = handler.extract("output.pptx")
    >>> print(result.watermark.content)

依赖:
    - python-pptx: PowerPoint 文件操作

Author: StealthMark Team
Date: 2026-04-28
"""

import os
import zipfile
import xml.etree.ElementTree as ET
from typing import Optional, Dict, Any
import logging

# 延迟导入
try:
    import pptx
except ImportError:
    pptx = None

from ..core.base import (
    BaseHandler, WatermarkData, WatermarkStatus,
    EmbedResult, ExtractResult, VerifyResult
)

# 模块日志
logger = logging.getLogger(__name__)


class PPTXHandler(BaseHandler):
    """
    PowerPoint演示文稿水印处理器
    
    使用隐藏标记方式嵌入水印。
    
    水印标记格式:
        hidden_<水印内容>_<原文>
    
    嵌入位置:
        每个幻灯片中所有文本框的文本前
    
    Attributes:
        无特殊属性，继承自 BaseHandler
    
    限制:
        - 水印标记在编辑时可见
        - 可通过"查找替换"删除
        - 不适合高安全性场景
    
    改进方向:
        - 使用隐藏形状（不可见矩形）
        - 使用备注页存储
        - 使用自定义 XML 部件
    """
    
    SUPPORTED_EXTENSIONS = ('.pptx',)
    HANDLER_NAME = "pptx"
    
    def embed(self, file_path: str, watermark: WatermarkData,
              output_path: str, **kwargs) -> EmbedResult:
        """
        嵌入水印到 PowerPoint
        
        嵌入流程:
        1. 打开演示文稿
        2. 遍历所有幻灯片
        3. 遍历每个幻灯片的所有形状
        4. 在文本框的文本前添加水印标记
        5. 保存演示文稿
        
        Args:
            file_path: 原始 PowerPoint 文件路径
            watermark: 水印数据对象
            output_path: 输出文件路径
            **kwargs: 额外参数（未使用）
        
        Returns:
            EmbedResult: 嵌入结果
        
        Note:
            水印会添加到所有文本框，确保即使删除部分幻灯片也能提取。
        """
        logger.info(f"PPTX embed: {file_path} -> {output_path}")
        
        # 检查依赖
        if pptx is None:
            logger.error("python-pptx not installed")
            return EmbedResult(
                status=WatermarkStatus.FAILED,
                message="需要安装python-pptx库: pip install python-pptx",
                file_path=file_path
            )
        
        # 验证文件
        error_result = self._validate_file(file_path)
        if error_result:
            return error_result
        
        try:
            # 打开演示文稿
            prs = pptx.Presentation(file_path)
            
            # 遍历所有幻灯片
            slide_count = 0
            for slide in prs.slides:
                slide_count += 1
                
                # 遍历所有形状
                for shape in slide.shapes:
                    if shape.has_text_frame:
                        # 在每个段落前添加水印标记
                        for paragraph in shape.text_frame.paragraphs:
                            for run in paragraph.runs:
                                # 添加水印标记
                                original_text = run.text
                                run.text = "hidden_" + watermark.content + "_" + original_text
            
            # 保存
            prs.save(output_path)
            
            logger.info(f"PPTX embed success: {slide_count} slides processed")
            return self._create_success_result(output_path)
            
        except Exception as e:
            logger.error(f"PPTX embed failed: {e}")
            return EmbedResult(
                status=WatermarkStatus.FAILED,
                message=f"嵌入失败: {str(e)}",
                file_path=file_path
            )
    
    def extract(self, file_path: str, **kwargs) -> ExtractResult:
        """
        从 PowerPoint 提取水印
        
        提取流程:
        1. 解压 pptx 文件（ZIP 格式）
        2. 遍历所有幻灯片 XML 文件
        3. 查找 "hidden_" 标记
        4. 提取标记后的水印内容
        
        Args:
            file_path: PowerPoint 文件路径
            **kwargs: 额外参数（未使用）
        
        Returns:
            ExtractResult: 提取结果
        
        Note:
            直接解析 XML 而非使用 python-pptx，更可靠。
            只需找到第一个水印标记即可。
        """
        logger.info(f"PPTX extract: {file_path}")
        
        # 检查依赖
        if pptx is None:
            logger.error("python-pptx not installed")
            return ExtractResult(
                status=WatermarkStatus.FAILED,
                message="需要安装python-pptx库",
                file_path=file_path
            )
        
        # 验证文件
        error_result = self._validate_file(file_path)
        if error_result:
            return ExtractResult(
                status=error_result.status,
                message=error_result.message,
                file_path=file_path
            )
        
        try:
            # 解压 pptx（ZIP 格式）
            with zipfile.ZipFile(file_path, 'r') as zf:
                files = zf.namelist()
                
                # 遍历所有幻灯片
                for fname in files:
                    if fname.startswith('ppt/slides/slide') and fname.endswith('.xml'):
                        with zf.open(fname) as f:
                            content = f.read().decode('utf-8')
                        
                        # 查找水印标记
                        if 'hidden_' in content:
                            # 提取水印内容
                            start = content.find('hidden_') + 7  # 跳过 "hidden_"
                            end = content.find('_', start)  # 找到结束标记
                            
                            if start > 6 and end > start:
                                watermark_text = content[start:end]
                                
                                logger.info(f"PPTX extract success: {watermark_text[:30]}...")
                                return ExtractResult(
                                    status=WatermarkStatus.SUCCESS,
                                    message="水印提取成功",
                                    file_path=file_path,
                                    watermark=WatermarkData(content=watermark_text)
                                )
            
            # 未找到水印
            logger.warning("No watermark found in PPTX")
            return ExtractResult(
                status=WatermarkStatus.EXTRACTION_FAILED,
                message="未找到水印",
                file_path=file_path
            )
            
        except Exception as e:
            logger.error(f"PPTX extract failed: {e}")
            return ExtractResult(
                status=WatermarkStatus.EXTRACTION_FAILED,
                message=f"提取失败: {str(e)}",
                file_path=file_path
            )
    
    def verify(self, file_path: str, original_watermark: WatermarkData,
               **kwargs) -> VerifyResult:
        """
        验证 PowerPoint 水印
        
        Args:
            file_path: PowerPoint 文件路径
            original_watermark: 原始水印数据
            **kwargs: 额外参数
        
        Returns:
            VerifyResult: 验证结果
        """
        logger.info(f"PPTX verify: {file_path}")
        
        # 提取水印
        extract_result = self.extract(file_path)
        
        # 提取失败
        if not extract_result.is_success or not extract_result.watermark:
            return VerifyResult(
                status=extract_result.status,
                is_valid=False,
                is_integrity_ok=False,
                match_score=0.0,
                message=f"提取失败: {extract_result.message}"
            )
        
        # 比对
        extracted = extract_result.watermark.content
        original = original_watermark.content
        
        is_match = extracted == original
        match_score = 1.0 if is_match else 0.0
        
        result = VerifyResult(
            status=WatermarkStatus.SUCCESS if is_match else WatermarkStatus.VERIFICATION_FAILED,
            is_valid=is_match,
            is_integrity_ok=True,
            match_score=match_score,
            message="验证通过" if is_match else "水印不匹配"
        )
        
        logger.info(f"PPTX verify result: valid={is_match}")
        return result


# 模块初始化日志
logger.info(f"{__name__} module loaded - PPTX handler ready")
